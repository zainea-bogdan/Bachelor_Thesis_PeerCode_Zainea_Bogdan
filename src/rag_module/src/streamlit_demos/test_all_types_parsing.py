import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import streamlit as st
import tempfile

# ── inline services (no import path dependency) ──────────────

from docx import Document as DocxDocument
from pptx import Presentation
import fitz


class DocxParsingService:
    PARAGRAPH_WINDOW = 200
    SECTION_WINDOW = 500

    def parse_docx(self, path: str) -> list[dict]:
        doc = DocxDocument(path)
        headings_found = []
        for paragraph in doc.paragraphs:
            if paragraph.style.name.startswith("Heading"):
                headings_found.append(paragraph)
        if len(headings_found) == 0:
            return self._chunk_by_paragraph(path)
        return self._chunk_by_headings(path)

    def _chunk_by_paragraph(self, path: str) -> list[dict]:
        doc = DocxDocument(path)
        filename = os.path.basename(path)
        chunks = []
        chunk_index = 0
        current_words = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if self._is_noise(text):
                continue
            words_in_this_paragraph = text.split()
            for word in words_in_this_paragraph:
                current_words.append(word)
            if len(current_words) >= self.PARAGRAPH_WINDOW:
                chunk_text = " ".join(current_words)
                chunks.append(self._build_chunk(chunk_text, filename, chunk_index, "paragraph_chunking"))
                chunk_index = chunk_index + 1
                current_words = []
        if len(current_words) > 0:
            chunk_text = " ".join(current_words)
            chunks.append(self._build_chunk(chunk_text, filename, chunk_index, "paragraph_chunking"))
        return chunks

    def _chunk_by_headings(self, path: str) -> list[dict]:
        doc = DocxDocument(path)
        filename = os.path.basename(path)
        chunks = []
        chunk_index = 0
        current_heading = None
        current_paragraphs = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if self._is_noise(text):
                continue
            if paragraph.style.name.startswith("Heading"):
                if len(current_paragraphs) > 0:
                    section_chunks = self._flush_section(current_heading, current_paragraphs, filename, chunk_index)
                    for chunk in section_chunks:
                        chunks.append(chunk)
                    chunk_index = chunk_index + len(section_chunks)
                current_heading = text
                current_paragraphs = []
            else:
                current_paragraphs.append(text)
        if len(current_paragraphs) > 0:
            section_chunks = self._flush_section(current_heading, current_paragraphs, filename, chunk_index)
            for chunk in section_chunks:
                chunks.append(chunk)
        return chunks

    def _flush_section(self, heading, paragraphs, filename, chunk_index_start):
        if heading is not None:
            full_text = heading + "\n\n" + " ".join(paragraphs)
        else:
            full_text = " ".join(paragraphs)
        word_count = len(full_text.split())
        if word_count <= self.SECTION_WINDOW:
            return [self._build_chunk(full_text, filename, chunk_index_start, "heading_chunking")]
        return self._split_into_windows(full_text, filename, chunk_index_start, "heading_chunking_split_oversized")

    def _split_into_windows(self, text, filename, chunk_index_start, strategy):
        all_words = text.split()
        chunks = []
        chunk_index = chunk_index_start
        position = 0
        while position < len(all_words):
            window_words = all_words[position: position + self.SECTION_WINDOW]
            chunks.append(self._build_chunk(" ".join(window_words), filename, chunk_index, strategy))
            chunk_index = chunk_index + 1
            position = position + self.SECTION_WINDOW
        return chunks

    def _build_chunk(self, text, filename, chunk_index, strategy):
        return {"text": text, "metadata": {"source_filename": filename, "document_type": "docx", "chunk_index": chunk_index, "strategy_used": strategy}}

    @staticmethod
    def _is_noise(text):
        if not text:
            return True
        if len(text.split()) < 3:
            return True
        return False


class PdfParsingService:
    SPARSE_THRESHOLD = 30
    BUFFER_CEILING = 100
    SECTION_WINDOW = 500

    def parse_pdf(self, path: str) -> list[dict]:
        filename = os.path.basename(path)
        chunks = []
        chunk_index = 0
        buffer_text = []
        buffer_words = 0
        try:
            doc = fitz.open(path)
        except Exception as e:
            raise Exception(f"Failed to open PDF {filename}: {str(e)}")
        for page_index in range(len(doc)):
            page = doc[page_index]
            text = page.get_text().strip()
            if self._is_noise(text):
                continue
            wc = len(text.split())
            if wc < self.SPARSE_THRESHOLD:
                if buffer_words >= self.BUFFER_CEILING:
                    flushed = self._flush_chunk(" ".join(buffer_text), filename, chunk_index, "sparse_page_chunking")
                    for chunk in flushed:
                        chunks.append(chunk)
                    chunk_index = chunk_index + len(flushed)
                    buffer_text = []
                    buffer_words = 0
                buffer_text.append(text)
                buffer_words = buffer_words + wc
            else:
                if len(buffer_text) > 0:
                    combined_text = " ".join(buffer_text) + " " + text
                    combined = self._flush_chunk(combined_text, filename, chunk_index, "combined_page_chunking")
                    for chunk in combined:
                        chunks.append(chunk)
                    chunk_index = chunk_index + len(combined)
                    buffer_text = []
                    buffer_words = 0
                else:
                    page_chunks = self._flush_chunk(text, filename, chunk_index, "page_chunking")
                    for chunk in page_chunks:
                        chunks.append(chunk)
                    chunk_index = chunk_index + len(page_chunks)
        doc.close()
        if len(buffer_text) > 0:
            remaining = self._flush_chunk(" ".join(buffer_text), filename, chunk_index, "sparse_page_chunking")
            for chunk in remaining:
                chunks.append(chunk)
        return chunks

    def _flush_chunk(self, text, filename, chunk_index_start, strategy):
        word_count = len(text.split())
        if word_count <= self.SECTION_WINDOW:
            return [self._build_chunk(text, filename, chunk_index_start, strategy)]
        return self._split_into_windows(text, filename, chunk_index_start, strategy)

    def _split_into_windows(self, text, filename, chunk_index_start, strategy):
        all_words = text.split()
        chunks = []
        chunk_index = chunk_index_start
        position = 0
        while position < len(all_words):
            window_words = all_words[position: position + self.SECTION_WINDOW]
            chunks.append(self._build_chunk(" ".join(window_words), filename, chunk_index, strategy + "_split"))
            chunk_index = chunk_index + 1
            position = position + self.SECTION_WINDOW
        return chunks

    def _build_chunk(self, text, filename, chunk_index, strategy):
        return {"text": text, "metadata": {"source_filename": filename, "document_type": "pdf", "chunk_index": chunk_index, "strategy_used": strategy}}

    @staticmethod
    def _is_noise(text):
        if not text:
            return True
        if len(text.split()) < 3:
            return True
        return False


class PptxParsingService:
    SECTION_WINDOW = 500

    def parse_pptx(self, path: str) -> list[dict]:
        filename = os.path.basename(path)
        chunks = []
        chunk_index = 0
        try:
            presentation = Presentation(path)
        except Exception as e:
            raise Exception(f"Failed to open PPTX {filename}: {str(e)}")
        for slide_index in range(len(presentation.slides)):
            slide = presentation.slides[slide_index]
            shape_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_text = shape.text.strip()
                    if len(shape_text) > 0:
                        shape_texts.append(shape_text)
            slide_number = slide_index + 1
            slide_text = f"Slide {slide_number}\n\n" + " ".join(shape_texts)
            if self._is_noise(slide_text):
                continue
            slide_chunks = self._flush_chunk(slide_text, filename, chunk_index)
            for chunk in slide_chunks:
                chunks.append(chunk)
            chunk_index = chunk_index + len(slide_chunks)
        return chunks

    def _flush_chunk(self, text, filename, chunk_index_start):
        word_count = len(text.split())
        if word_count <= self.SECTION_WINDOW:
            return [self._build_chunk(text, filename, chunk_index_start, "slide_chunking")]
        return self._split_into_windows(text, filename, chunk_index_start)

    def _split_into_windows(self, text, filename, chunk_index_start):
        all_words = text.split()
        chunks = []
        chunk_index = chunk_index_start
        position = 0
        while position < len(all_words):
            window_words = all_words[position: position + self.SECTION_WINDOW]
            chunks.append(self._build_chunk(" ".join(window_words), filename, chunk_index, "slide_chunking_split"))
            chunk_index = chunk_index + 1
            position = position + self.SECTION_WINDOW
        return chunks

    def _build_chunk(self, text, filename, chunk_index, strategy):
        return {"text": text, "metadata": {"source_filename": filename, "document_type": "pptx", "chunk_index": chunk_index, "strategy_used": strategy}}

    @staticmethod
    def _is_noise(text):
        if not text:
            return True
        if len(text.split()) < 3:
            return True
        return False


# ── app ──────────────────────────────────────────────────────

st.set_page_config(page_title="Parsing test harness — all types", layout="wide")

docx_service = DocxParsingService()
pdf_service = PdfParsingService()
pptx_service = PptxParsingService()

# ── sidebar — simulates node js + teacher inputs ─────────────

st.sidebar.title("Simulated teacher inputs")
st.sidebar.caption("Simulates what Node.js would send to FastAPI")
st.sidebar.markdown("---")

uploaded_file = st.sidebar.file_uploader(
    "Upload document",
    type=["pdf", "docx", "pptx"]
)

st.sidebar.markdown("---")
st.sidebar.markdown("**Simulated metadata (from Node.js)**")

course_id = st.sidebar.text_input("course_id", value="CS101")
teacher_id = st.sidebar.text_input("teacher_id", value="T001")
university_year = st.sidebar.selectbox(
    "university_year",
    ["2023-2024", "2024-2025", "2025-2026"],
    index=1
)

st.sidebar.markdown("---")
st.sidebar.markdown("**Display options**")
show_full_text = st.sidebar.checkbox("Show full chunk text", value=False)
flag_below = st.sidebar.slider("Flag chunks below (words)", 0, 100, 20)

# ── main ─────────────────────────────────────────────────────

st.title("Parsing test harness — all document types")
st.caption("Simulates the full parsing pipeline: teacher uploads → Node.js forwards → FastAPI parses → chunks ready for embedding")

if not uploaded_file:
    st.info("Upload a PDF, DOCX, or PPTX from the sidebar to begin.")

    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("**DOCX strategy**")
        st.markdown("""
        - Headings detected → heading chunking
        - No headings → paragraph chunking
        - Ceiling: 500 words per chunk
        - Window: 200 words per paragraph chunk
        """)
    with col2:
        st.markdown("**PDF strategy**")
        st.markdown("""
        - Sparse pages (&lt;30w) → buffer
        - Buffer ceiling: 100 words
        - Substantial page + buffer → combine
        - Ceiling: 500 words per chunk
        """)
    with col3:
        st.markdown("**PPTX strategy**")
        st.markdown("""
        - 1 slide = 1 chunk
        - Slide number embedded in text
        - Ceiling: 500 words per chunk
        - No buffer needed
        """)
    st.stop()

# ── detect file type ─────────────────────────────────────────

file_ext = uploaded_file.name.split(".")[-1].lower()
file_size_kb = round(uploaded_file.size / 1024, 1)

# ── save to temp file ─────────────────────────────────────────

with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as tmp:
    tmp.write(uploaded_file.read())
    tmp_path = tmp.name

# ── document overview ─────────────────────────────────────────

st.subheader("Document overview")

col1, col2, col3, col4 = st.columns(4)
col1.metric("Filename", uploaded_file.name)
col2.metric("File type", file_ext.upper())
col3.metric("File size", f"{file_size_kb} KB")
col4.metric("Simulated course", course_id)

# ── show document-specific info ───────────────────────────────

if file_ext == "docx":
    doc = DocxDocument(tmp_path)
    headings = [p for p in doc.paragraphs if p.style.name.startswith("Heading")]
    total_words = sum(len(p.text.split()) for p in doc.paragraphs if p.text.strip())
    total_paragraphs = len([p for p in doc.paragraphs if p.text.strip()])

    m1, m2, m3, m4 = st.columns(4)
    m1.metric("Total paragraphs", total_paragraphs)
    m2.metric("Headings found", len(headings))
    m3.metric("Total words", total_words)
    m4.metric("Strategy", "heading chunking" if len(headings) > 0 else "paragraph chunking")

    if len(headings) > 0:
        st.success(f"Headings detected — heading-based chunking will be used")
        with st.expander("Heading structure preview"):
            for h in headings:
                level = h.style.name.replace("Heading ", "")
                indent = "&nbsp;" * (int(level) * 4) if level.isdigit() else ""
                st.markdown(f"{indent}**{h.style.name}** — {h.text}", unsafe_allow_html=True)
    else:
        st.warning("No headings detected — paragraph chunking (200-word window) will be used")

elif file_ext == "pdf":
    raw_doc = fitz.open(tmp_path)
    total_pages = len(raw_doc)
    raw_doc.close()

    m1, m2 = st.columns(2)
    m1.metric("Total pages", total_pages)
    m2.metric("Strategy", "adaptive buffer (sparse/substantial detection)")

elif file_ext == "pptx":
    prs = Presentation(tmp_path)
    total_slides = len(prs.slides)

    m1, m2 = st.columns(2)
    m1.metric("Total slides", total_slides)
    m2.metric("Strategy", "1 slide = 1 chunk")

st.markdown("---")

# ── run parsing ───────────────────────────────────────────────

run = st.button("Run parsing", type="primary")

if not run:
    os.unlink(tmp_path)
    st.stop()

with st.spinner("Parsing document..."):
    try:
        if file_ext == "docx":
            chunks = docx_service.parse_docx(tmp_path)
        elif file_ext == "pdf":
            chunks = pdf_service.parse_pdf(tmp_path)
        elif file_ext == "pptx":
            chunks = pptx_service.parse_pptx(tmp_path)
        else:
            st.error(f"Unsupported file type: {file_ext}")
            os.unlink(tmp_path)
            st.stop()
    except Exception as e:
        st.error(f"Parsing failed: {e}")
        os.unlink(tmp_path)
        st.stop()

os.unlink(tmp_path)

if not chunks:
    st.error("No chunks produced — document may be empty or image-only.")
    st.stop()

# ── results summary ───────────────────────────────────────────

st.subheader("Parsing results")

word_counts = [len(c["text"].split()) for c in chunks]
total_words_chunked = sum(word_counts)
avg_words = round(total_words_chunked / len(chunks), 1)
sparse_chunks = [c for c in chunks if len(c["text"].split()) < flag_below]

m1, m2, m3, m4, m5 = st.columns(5)
m1.metric("Total chunks", len(chunks))
m2.metric("Total words chunked", total_words_chunked)
m3.metric("Avg words/chunk", avg_words)
m4.metric("Min words/chunk", min(word_counts))
m5.metric("Max words/chunk", max(word_counts))

if len(sparse_chunks) > 0:
    st.warning(f"{len(sparse_chunks)} chunk(s) below {flag_below} words — may be too sparse for embedding")
else:
    st.success("All chunks above the sparse threshold")

# ── strategy breakdown ────────────────────────────────────────

st.markdown("**Strategy breakdown:**")
strategy_counts = {}
for c in chunks:
    s = c["metadata"]["strategy_used"]
    strategy_counts[s] = strategy_counts.get(s, 0) + 1

for strategy, count in strategy_counts.items():
    pct = round(count / len(chunks) * 100, 1)
    st.write(f"- `{strategy}`: {count} chunks ({pct}%)")

# ── word count distribution ───────────────────────────────────

st.markdown("---")
st.subheader("Word count distribution")

buckets = {"< 50": 0, "50–100": 0, "100–200": 0, "200–400": 0, "400–500": 0, "> 500": 0}
for wc in word_counts:
    if wc < 50:
        buckets["< 50"] += 1
    elif wc < 100:
        buckets["50–100"] += 1
    elif wc < 200:
        buckets["100–200"] += 1
    elif wc < 400:
        buckets["200–400"] += 1
    elif wc <= 500:
        buckets["400–500"] += 1
    else:
        buckets["> 500"] += 1

max_count = max(buckets.values()) if buckets.values() else 1
for label, count in buckets.items():
    bar_width = int((count / max_count) * 100) if max_count > 0 else 0
    col_label, col_bar, col_count = st.columns([2, 5, 1])
    with col_label:
        st.markdown(f"<small>{label} words</small>", unsafe_allow_html=True)
    with col_bar:
        color = "#E24B4A" if label == "> 500" else "#2E75B6"
        if count > 0:
            st.markdown(
                f'<div style="background:{color};height:20px;width:{bar_width}%;border-radius:3px;margin-top:4px"></div>',
                unsafe_allow_html=True
            )
    with col_count:
        st.markdown(f"<small>**{count}**</small>", unsafe_allow_html=True)

if buckets["> 500"] > 0:
    st.error(f"{buckets['> 500']} chunk(s) exceed 500 words — Liu et al. 2023 ceiling breached. Check splitting logic.")

# ── simulated chromadb payload preview ───────────────────────

st.markdown("---")
st.subheader("Simulated ChromaDB payload")
st.caption("This is what IngestionService will send to ChromaDB — metadata enriched with course and teacher context from Node.js")

sample_chunk = chunks[0]
simulated_payload = {
    "id": f"{course_id}_{uploaded_file.name}_chunk_0",
    "text": sample_chunk["text"][:200] + "..." if len(sample_chunk["text"]) > 200 else sample_chunk["text"],
    "metadata": {
        **sample_chunk["metadata"],
        "course_id": course_id,
        "teacher_id": teacher_id,
        "university_year": university_year,
        "upload_timestamp": "2025-05-09T00:00:00Z"
    }
}
st.json(simulated_payload)

# ── chunk inspector ───────────────────────────────────────────

st.markdown("---")
st.subheader("Chunk inspector")

search_term = st.text_input("Search chunks by content", placeholder="type to filter...")

filtered_chunks = chunks
if search_term:
    filtered_chunks = [c for c in chunks if search_term.lower() in c["text"].lower()]
    st.caption(f"{len(filtered_chunks)} of {len(chunks)} chunks match '{search_term}'")

for i, chunk in enumerate(filtered_chunks):
    wc = len(chunk["text"].split())
    is_sparse = wc < flag_below
    is_oversized = wc > 500

    label = f"Chunk {chunk['metadata']['chunk_index']} — {wc} words — `{chunk['metadata']['strategy_used']}`"
    if is_sparse:
        label = "⚠ " + label + " (sparse)"
    if is_oversized:
        label = "🔴 " + label + " (oversized)"

    with st.expander(label, expanded=False):
        tab1, tab2 = st.tabs(["Text", "Metadata"])

        with tab1:
            text_to_show = chunk["text"] if show_full_text else chunk["text"][:600]
            if not show_full_text and len(chunk["text"]) > 600:
                text_to_show += f"\n\n[...{len(chunk['text']) - 600} more characters — enable 'Show full chunk text']"
            st.text(text_to_show)

        with tab2:
            enriched_metadata = {
                **chunk["metadata"],
                "course_id": course_id,
                "teacher_id": teacher_id,
                "university_year": university_year
            }
            st.json(enriched_metadata)