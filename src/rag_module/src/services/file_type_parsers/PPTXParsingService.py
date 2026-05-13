import os
from pptx import Presentation
from src.services.file_type_parsers.BaseParsingService import BaseParsingService


class PptxParsingService(BaseParsingService):

    # slides below this word count are considered sparse
    # titles, labels, short bullets — too weak alone for embedding
    SPARSE_THRESHOLD = 30

    def _get_document_type(self) -> str:
        return "pptx"

    def parse(self, path: str) -> list[dict]:
        filename = os.path.basename(path)

        try:
            presentation = Presentation(path)
        except Exception as e:
            raise Exception(f"Failed to open PPTX {filename}: {str(e)}")

        chunks = []
        chunk_index = 0

        # buffer for sparse slides — short slides wait
        # for a substantial slide to attach to
        buffer_texts = []
        buffer_words = 0

        for slide_index in range(len(presentation.slides)):
            slide = presentation.slides[slide_index]

            # collect all shape texts on this slide
            shape_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if not self._is_noise(text):
                        shape_texts.append(text)

            # skip slide if no meaningful text
            if len(shape_texts) == 0:
                continue

            # build slide text
            slide_text = " ".join(shape_texts)
            slide_words = len(slide_text.split())

            if slide_words < self.SPARSE_THRESHOLD:
                # sparse slide — buffer it
                # title or label slide waits for explanation
                buffer_texts.append(slide_text)
                buffer_words = buffer_words + slide_words

            else:
                # substantial slide
                if len(buffer_texts) > 0:
                    # combine buffer + this slide → one chunk
                    # title travels with its explanation
                    combined_text = " ".join(buffer_texts) + " " + slide_text
                    chunks.append(self._build_chunk(
                        text=combined_text,
                        filename=filename,
                        chunk_index=chunk_index,
                        strategy="slide_chunking"
                    ))
                    chunk_index = chunk_index + 1
                    buffer_texts = []
                    buffer_words = 0

                else:
                    # no buffer — slide stands alone
                    chunks.append(self._build_chunk(
                        text=slide_text,
                        filename=filename,
                        chunk_index=chunk_index,
                        strategy="slide_chunking"
                    ))
                    chunk_index = chunk_index + 1

        # flush remaining buffer at end of presentation
        if len(buffer_texts) > 0:
            remaining_text = " ".join(buffer_texts)
            chunks.append(self._build_chunk(
                text=remaining_text,
                filename=filename,
                chunk_index=chunk_index,
                strategy="slide_chunking"
            ))

        return chunks