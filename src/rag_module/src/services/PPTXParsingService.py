import os
from pptx import Presentation


class PptxParsingService:

    # maximum words per chunk
    # justified by: Lost in the Middle — Liu et al. 2023
    SECTION_WINDOW = 500

    def parse_pptx(self, path: str) -> list[dict]:
        filename = os.path.basename(path)
        chunks = []
        chunk_index = 0

        try:
            presentation = Presentation(path)
        except Exception as e:
            raise Exception(f"Failed to open PPTX {filename}: {str(e)}")

        # iterate every slide in the presentation
        for slide_index in range(len(presentation.slides)):
            slide = presentation.slides[slide_index]

            # collect text from every shape on this slide
            shape_texts = []
            for shape in slide.shapes:

                # only process shapes that contain a text frame
                if shape.has_text_frame:
                    shape_text = shape.text.strip()

                    # skip empty shapes
                    if len(shape_text) > 0:
                        shape_texts.append(shape_text)

            # build the slide text with slide number as prefix
            # same pattern as DOCX heading chunking —
            # the slide number travels with its content
            slide_number = slide_index + 1
            slide_text = f"Slide {slide_number}\n\n" + " ".join(shape_texts)

            # skip slide if it produced no meaningful text
            if self._is_noise(slide_text):
                continue

            # flush the slide text as one or more chunks
            slide_chunks = self._flush_chunk(
                text=slide_text,
                filename=filename,
                chunk_index_start=chunk_index
            )

            for chunk in slide_chunks:
                chunks.append(chunk)

            chunk_index = chunk_index + len(slide_chunks)

        return chunks

    def _flush_chunk(
        self,
        text: str,
        filename: str,
        chunk_index_start: int
    ) -> list[dict]:
        # check word count against ceiling
        word_count = len(text.split())

        # fits within limit — emit as one chunk
        if word_count <= self.SECTION_WINDOW:
            return [self._build_chunk(
                text=text,
                filename=filename,
                chunk_index=chunk_index_start,
                strategy="slide_chunking"
            )]

        # too large — split into 500-word windows
        return self._split_into_windows(
            text=text,
            filename=filename,
            chunk_index_start=chunk_index_start
        )

    def _split_into_windows(
        self,
        text: str,
        filename: str,
        chunk_index_start: int
    ) -> list[dict]:
        # split the text into individual words
        all_words = text.split()

        chunks = []
        chunk_index = chunk_index_start

        # current position in the word list
        position = 0

        # slide a 500-word window across the words
        while position < len(all_words):

            # take the next 500 words
            # if fewer than 500 remain, take whatever is left
            window_words = all_words[position: position + self.SECTION_WINDOW]

            # join back into a string
            chunk_text = " ".join(window_words)

            # build and store this sub-chunk
            new_chunk = self._build_chunk(
                text=chunk_text,
                filename=filename,
                chunk_index=chunk_index,
                strategy="slide_chunking_split"
            )
            chunks.append(new_chunk)

            # move index and position forward
            chunk_index = chunk_index + 1
            position = position + self.SECTION_WINDOW

        return chunks

    def _build_chunk(
        self,
        text: str,
        filename: str,
        chunk_index: int,
        strategy: str
    ) -> dict:
        chunk = {
            "text": text,
            "metadata": {
                "source_filename": filename,
                "document_type": "pptx",
                "chunk_index": chunk_index,
                "strategy_used": strategy
            }
        }
        return chunk

    @staticmethod
    def _is_noise(text: str) -> bool:
        if not text:
            return True
        words = text.split()
        if len(words) < 3:
            return True
        return False